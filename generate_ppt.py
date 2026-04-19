import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title_text, subtitle_text, logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Styling
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(41, 128, 185) # Blue theme

    if logo_path and os.path.exists(logo_path):
        left = Inches(4.25)
        top = Inches(0.5)
        height = Inches(1.5)
        slide.shapes.add_picture(logo_path, left, top, height=height)

def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title.text = title_text
    tf = body_shape.text_frame
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            if point.startswith("  -") or point.startswith("-"):
                p.level = 1

def create_pitch_deck():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, "StudyCompass", "AI Student OS & Next-Gen Loan Funnel", logo_path="logo.png")

    # Slide 2: The Problem
    add_content_slide(prs, "The Problem", [
        "Traditional study abroad platforms are broken:",
        "  - Boring lead forms cause massive student drop-off.",
        "  - Advice is generic and lacks personalization.",
        "  - High friction to financial products (loans).",
        "  - Zero incentive for students to return after visiting once."
    ])

    # Slide 3: The Solution
    add_content_slide(prs, "The Solution: StudyCompass", [
        "A highly addictive, \"all-in-one\" OS for students targeting study abroad.",
        "  - Replaces boring forms with a dynamic 'Profile Score'.",
        "  - Provides free, high-value AI generators (SOPs, Timelines).",
        "  - Acts as a 'Trojan Horse' to capture high-intent financial data.",
        "  - Monetizes via a seamless, AI-guided student loan conversion funnel."
    ])

    # Slide 4: Key Feature 1 - "Weaponized" Loan Funnel
    add_content_slide(prs, "Feature 1: The Loan Funnel", [
        "The primary business driver connecting students to NBFCs/Banks:",
        "  - Calculates exact 'Loan Readiness' percentage.",
        "  - Estimates financial shortfall dynamically.",
        "  - Recommends the perfect loan product (e.g., Prodigy Finance).",
        "  - Uses the Claude AI Advisor to explain complex terminology instantly."
    ])

    # Slide 5: Key Feature 2 - High-Intent Engagement Engine
    add_content_slide(prs, "Feature 2: Engagement Engine & AI Copilot", [
        "Keeps the student engaged long enough to convert:",
        "  - Continuous 'Next Best Action' component pinned to the sidebar.",
        "  - Visual tracking of Profile Strength (gamification).",
        "  - Proactive AI Mentor that remembers previous states.",
        "  - Example: 'You mentioned ₹20L budget — I suggest exploring Germany...'"
    ])

    # Slide 6: Technology Stack
    add_content_slide(prs, "Under the Hood (Tech Stack)", [
        "- Frontend & State Management: Streamlit (Rapid, Python-native UI).",
        "- Intelligence Core: Anthropic Claude 3.5 Sonnet.",
        "- Why Claude 3.5?: Superior contextual memory and proactive suggestion modeling.",
        "- Architecture: Modular, lightweight, highly scalable Python backend."
    ])

    # Slide 7: Business Value
    add_content_slide(prs, "Why This Wins (Business Value)", [
        "- Extreme Explainability: Judges & clients love the 'Explain My Profile' AI transparency.",
        "- Real-world Usability: Export features (PDF plans) anchor it as a real product.",
        "- Unmatched Conversion rates: Warm hand-offs to partner NBFCs.",
        "- Scaleable: Replaces human counseling cost with instant API calls."
    ])

    prs.save('StudyCompass_Pitch.pptx')
    print("Presentation generated successfully: StudyCompass_Pitch.pptx")

if __name__ == '__main__':
    create_pitch_deck()
